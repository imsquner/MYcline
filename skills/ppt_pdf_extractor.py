# -*- coding: utf-8 -*-
"""
[技能] PPT/PDF 内容提取器
功能：从 .ppt/.pptx/.ppsx 和 .pdf 文件中提取文字内容
用法：python ppt_pdf_extractor.py <源目录> <输出目录>
依赖：python-pptx, PyMuPDF, pdfplumber
作者：Cline for MYcline
"""
import os, sys, re, zipfile, xml.etree.ElementTree as ET
sys.stdout.reconfigure(encoding='utf-8') if hasattr(sys.stdout, 'reconfigure') else None

def extract_text_from_slideshow(path):
    """从 slideshow 格式 (.ppt 实质是 .ppsx) 提取文字"""
    texts = []
    try:
        with zipfile.ZipFile(path, 'r') as z:
            slide_files = [f for f in z.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
            slide_files.sort(key=lambda x: int(re.search(r'slide(\d+)', x).group(1)))
            for sfile in slide_files:
                content = z.read(sfile)
                root = ET.fromstring(content)
                for t_elem in root.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}t'):
                    if t_elem.text and t_elem.text.strip():
                        texts.append(t_elem.text.strip())
        seen = set()
        result = []
        for t in texts:
            if t not in seen and len(t) > 3:
                seen.add(t)
                result.append(t)
        return result
    except Exception as e:
        print(f"    Slideshow提取失败: {e}")
        return None

def extract_ppt_text(path):
    """从 PPT/PPTX 文件提取文字"""
    from pptx import Presentation
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    t = shape.text.strip().replace("\n", " ")
                    if len(t) > 3:
                        texts.append(t)
        return list(dict.fromkeys(texts))
    except Exception:
        return extract_text_from_slideshow(path)

def extract_pdf_text(path):
    """从 PDF 提取文字"""
    import fitz
    doc = fitz.open(path)
    all_text = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text("text").strip()
        if text and len(text) > 10:
            all_text.append(f"--- Page {page_num+1} ---\n{text}")
    doc.close()
    if not all_text:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t and t.strip():
                    all_text.append(t.strip())
    return all_text

def process_directory(src_dir, out_dir):
    """处理目录下所有 PPT 和 PDF 文件"""
    os.makedirs(out_dir, exist_ok=True)
    files = [f for f in os.listdir(src_dir) if f.endswith(('.ppt','.pptx','.pps','.ppsx','.pdf'))]
    for fname in sorted(files):
        path = os.path.join(src_dir, fname)
        base = os.path.splitext(fname)[0]
        ext = os.path.splitext(fname)[1].lower()
        print(f"\n[{ext.upper()}] {fname}")
        if ext == '.pdf':
            texts = extract_pdf_text(path)
        else:
            texts = extract_ppt_text(path)
        if texts and len(texts) > 0:
            outpath = os.path.join(out_dir, f"{base}.md")
            with open(outpath, "w", encoding="utf-8") as f:
                f.write("\n\n".join(texts))
            print(f"  -> OK: {len(texts)} segments -> {outpath}")
        else:
            print(f"  -> 跳过: 无法提取内容")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("用法: python ppt_pdf_extractor.py <源目录> <输出目录>")
        sys.exit(1)
    process_directory(sys.argv[1], sys.argv[2])
